import logging
from pathlib import Path
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import io
from PIL import Image

log = logging.getLogger(__name__)

def extract_from_pdf(file_path: Path, output_dir: Path) -> dict:
    """從 PDF 檔案中提取所有文字和圖片。"""
    text_content = ""
    image_paths = []
    try:
        pdf_document = fitz.open(file_path)
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            text_content += page.get_text() + "\n"

            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]

                image_filename = output_dir / f"{file_path.stem}_page{page_num+1}_img{img_index}.{image_ext}"
                with open(image_filename, "wb") as img_file:
                    img_file.write(image_bytes)
                image_paths.append(image_filename)
        log.info(f"從 PDF '{file_path.name}' 中成功提取 {len(image_paths)} 張圖片和 {len(text_content)} 字元。")
    except Exception as e:
        log.error(f"從 PDF '{file_path.name}' 提取內容時發生錯誤: {e}", exc_info=True)
    return {"text": text_content.strip(), "image_paths": image_paths}

def extract_from_docx(file_path: Path, output_dir: Path) -> dict:
    """從 DOCX 檔案中提取所有文字和圖片。"""
    text_content = ""
    image_paths = []
    try:
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            text_content += para.text + "\n"

        for i, rel in enumerate(doc.part.rels.values()):
            if "image" in rel.target_ref:
                image_data = rel.target_part.blob
                try:
                    image = Image.open(io.BytesIO(image_data))
                    ext = image.format.lower()
                except Exception:
                    ext = 'png'
                image_filename = output_dir / f"{file_path.stem}_img{i}.{ext}"
                with open(image_filename, "wb") as img_file:
                    img_file.write(image_data)
                image_paths.append(image_filename)
        log.info(f"從 DOCX '{file_path.name}' 中成功提取 {len(image_paths)} 張圖片和 {len(text_content)} 字元。")
    except Exception as e:
        log.error(f"從 DOCX '{file_path.name}' 提取內容時發生錯誤: {e}", exc_info=True)
    return {"text": text_content.strip(), "image_paths": image_paths}

def extract_from_pptx(file_path: Path, output_dir: Path) -> dict:
    """從 PPTX 檔案中提取所有文字和圖片。"""
    text_content = ""
    image_paths = []
    try:
        prs = Presentation(file_path)
        img_index = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image = shape.image
                    image_bytes = image.blob
                    ext = image.ext
                    image_filename = output_dir / f"{file_path.stem}_img{img_index}.{ext}"
                    with open(image_filename, "wb") as img_file:
                        img_file.write(image_bytes)
                    image_paths.append(image_filename)
                    img_index += 1
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_content += run.text + " "
                    text_content += "\n"
        log.info(f"從 PPTX '{file_path.name}' 中成功提取 {len(image_paths)} 張圖片和 {len(text_content)} 字元。")
    except Exception as e:
        log.error(f"從 PPTX '{file_path.name}' 提取內容時發生錯誤: {e}", exc_info=True)
    return {"text": text_content.strip(), "image_paths": image_paths}

def extract_content(file_path_str: str, image_output_dir_str: str) -> dict | None:
    """
    一個主函式，根據副檔名分派任務給對應的提取器。
    現在會同時提取文字與圖片。

    :param file_path_str: 來源檔案的完整路徑字串。
    :param image_output_dir_str: 儲存提取出的圖片的目錄路徑字串。
    :return: 一個包含 'text' 和 'image_paths' 的字典，或在失敗時回傳 None。
    """
    file_path = Path(file_path_str)
    output_dir = Path(image_output_dir_str)
    output_dir.mkdir(parents=True, exist_ok=True)

    if not file_path.is_file():
        log.error(f"檔案不存在於: {file_path}")
        return None

    ext = file_path.suffix.lower()
    content_data = {"text": "", "image_paths": []}

    if ext == '.pdf':
        content_data = extract_from_pdf(file_path, output_dir)
    elif ext == '.docx':
        content_data = extract_from_docx(file_path, output_dir)
    elif ext == '.pptx':
        content_data = extract_from_pptx(file_path, output_dir)
    else:
        log.warning(f"不支援的檔案類型: {ext}。跳過內容提取。")
        return {"text": "", "image_paths": []}

    # 確保圖片路徑是字串格式
    content_data["image_paths"] = [str(p) for p in content_data["image_paths"]]
    return content_data
